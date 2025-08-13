import os
import fitz  # PyMuPDF
import docx
import openpyxl
from pptx import Presentation
import chardet
from typing import List, Dict, Any, Optional, Tuple
from datetime import datetime, timezone
import hashlib
import uuid
import logging
from qdrant_client import QdrantClient, models
from openai import OpenAI
from dotenv import load_dotenv

load_dotenv()

class DocumentProcessor:
    """Processes uploaded documents and stores them in Qdrant vector database"""
    
    def __init__(self):
        # Initialize clients
        self.qdrant_client = QdrantClient(
            url=os.getenv("QDRANT_URL"),
            api_key=os.getenv("QDRANT_API_KEY")
        )
        self.openai_client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
        self.collection_name = "vaccine_guidelines_2"
        
        # Supported file types
        self.supported_types = {
            '.pdf': self._extract_pdf_text,
            '.docx': self._extract_docx_text,
            '.doc': self._extract_docx_text,  # Will attempt docx processing
            '.pptx': self._extract_pptx_text,
            '.xlsx': self._extract_xlsx_text,
            '.xls': self._extract_xlsx_text,
            '.txt': self._extract_txt_text,
        }
        
        # Setup logging
        self.logger = logging.getLogger(__name__)
    
    def _extract_pdf_text(self, file_path: str) -> str:
        """Extract text from PDF file"""
        try:
            doc = fitz.open(file_path)
            text = ""
            for page in doc:
                text += page.get_text()
            doc.close()
            return text
        except Exception as e:
            self.logger.error(f"Error extracting PDF text: {e}")
            return ""
    
    def _extract_docx_text(self, file_path: str) -> str:
        """Extract text from DOCX file"""
        try:
            doc = docx.Document(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text
        except Exception as e:
            self.logger.error(f"Error extracting DOCX text: {e}")
            return ""
    
    def _extract_pptx_text(self, file_path: str) -> str:
        """Extract text from PPTX file"""
        try:
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except Exception as e:
            self.logger.error(f"Error extracting PPTX text: {e}")
            return ""
    
    def _extract_xlsx_text(self, file_path: str) -> str:
        """Extract text from XLSX file"""
        try:
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            text = ""
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                for row in sheet.iter_rows(values_only=True):
                    row_text = " ".join([str(cell) for cell in row if cell is not None])
                    if row_text.strip():
                        text += row_text + "\n"
            return text
        except Exception as e:
            self.logger.error(f"Error extracting XLSX text: {e}")
            return ""
    
    def _extract_txt_text(self, file_path: str) -> str:
        """Extract text from TXT file"""
        try:
            # Detect encoding
            with open(file_path, 'rb') as f:
                raw_data = f.read()
                encoding = chardet.detect(raw_data)['encoding']
            
            # Read with detected encoding
            with open(file_path, 'r', encoding=encoding or 'utf-8') as f:
                return f.read()
        except Exception as e:
            self.logger.error(f"Error extracting TXT text: {e}")
            return ""
    
    def extract_text_from_file(self, file_path: str, filename: str) -> str:
        """Extract text from supported file types"""
        file_ext = os.path.splitext(filename)[1].lower()
        
        if file_ext not in self.supported_types:
            raise ValueError(f"Unsupported file type: {file_ext}")
        
        extractor = self.supported_types[file_ext]
        text = extractor(file_path)
        
        if not text.strip():
            raise ValueError(f"No text content extracted from file: {filename}")
        
        return text
    
    def chunk_text(self, text: str, chunk_size: int = 1000, overlap: int = 200) -> List[str]:
        """Split text into overlapping chunks"""
        if len(text) <= chunk_size:
            return [text]
        
        chunks = []
        start = 0
        
        while start < len(text):
            end = start + chunk_size
            
            # Try to break at sentence boundary
            if end < len(text):
                # Look for sentence endings
                for i in range(end, max(start + chunk_size - 100, start), -1):
                    if text[i] in '.!?':
                        end = i + 1
                        break
                # If no sentence boundary found, break at word boundary
                else:
                    for i in range(end, max(start + chunk_size - 50, start), -1):
                        if text[i].isspace():
                            end = i
                            break
            
            chunk = text[start:end].strip()
            if chunk:
                chunks.append(chunk)
            
            start = end - overlap
            if start >= len(text):
                break
        
        return chunks
    
    def generate_embeddings(self, texts: List[str]) -> List[List[float]]:
        """Generate embeddings for text chunks"""
        try:
            response = self.openai_client.embeddings.create(
                model="text-embedding-3-small",
                input=texts
            )
            return [embedding.embedding for embedding in response.data]
        except Exception as e:
            self.logger.error(f"Error generating embeddings: {e}")
            raise
    
    def store_in_qdrant(self, chunks: List[str], embeddings: List[List[float]], 
                       metadata: Dict[str, Any]) -> str:
        """Store document chunks and embeddings in Qdrant"""
        try:
            # Generate document ID
            doc_id = str(uuid.uuid4())
            
            # Prepare points for insertion
            points = []
            for i, (chunk, embedding) in enumerate(zip(chunks, embeddings)):
                point_id = str(uuid.uuid4())
                payload = {
                    'text': chunk,
                    'doc_id': doc_id,
                    'chunk_id': i + 1,
                    'total_chunks': len(chunks),
                    'filename': metadata.get('filename', 'unknown'),
                    'source': metadata.get('source', 'user_upload'),
                    'source_type': metadata.get('source_type', 'document'),
                    'title': metadata.get('title', metadata.get('filename', 'unknown')),
                    'publication_date': metadata.get('publication_date', datetime.now(timezone.utc).isoformat()),
                    'upload_timestamp': datetime.now(timezone.utc).isoformat(),
                    'file_size': metadata.get('file_size', 0),
                    'file_type': metadata.get('file_type', 'unknown')
                }
                
                points.append({
                    'id': point_id,
                    'vector': embedding,
                    'payload': payload
                })
            
            # Insert points into Qdrant
            self.qdrant_client.upsert(
                collection_name=self.collection_name,
                points=points
            )
            
            self.logger.info(f"Successfully stored {len(chunks)} chunks for document {doc_id}")
            return doc_id
            
        except Exception as e:
            self.logger.error(f"Error storing in Qdrant: {e}")
            raise
    
    def process_uploaded_file(self, file_path: str, filename: str, 
                            metadata: Optional[Dict[str, Any]] = None) -> Dict[str, Any]:
        """Complete pipeline for processing uploaded file"""
        try:
            # Extract basic metadata
            file_stats = os.stat(file_path)
            file_size = file_stats.st_size
            file_type = os.path.splitext(filename)[1].lower()
            
            # Prepare metadata
            if metadata is None:
                metadata = {}
            
            metadata.update({
                'filename': filename,
                'file_size': file_size,
                'file_type': file_type,
                'source_type': 'user_upload',  # Mark as user upload
                'upload_timestamp': datetime.now(timezone.utc).isoformat()
            })
            
            # Extract text
            self.logger.info(f"Extracting text from {filename}")
            text = self.extract_text_from_file(file_path, filename)
            
            # Chunk text
            self.logger.info(f"Chunking text from {filename}")
            chunks = self.chunk_text(text)
            
            # Generate embeddings
            self.logger.info(f"Generating embeddings for {len(chunks)} chunks")
            embeddings = self.generate_embeddings(chunks)
            
            # Store in Qdrant
            self.logger.info(f"Storing document in Qdrant")
            doc_id = self.store_in_qdrant(chunks, embeddings, metadata)
            
            # Clean up temporary file
            if os.path.exists(file_path):
                os.remove(file_path)
            
            return {
                'success': True,
                'doc_id': doc_id,
                'filename': filename,
                'chunks_count': len(chunks),
                'file_size': file_size,
                'file_type': file_type,
                'upload_timestamp': metadata['upload_timestamp']
            }
            
        except Exception as e:
            # Clean up temporary file on error
            if os.path.exists(file_path):
                os.remove(file_path)
            
            self.logger.error(f"Error processing file {filename}: {e}")
            return {
                'success': False,
                'error': str(e),
                'filename': filename
            }
    
    def get_supported_file_types(self) -> List[str]:
        """Get list of supported file extensions"""
        return list(self.supported_types.keys())
    
    def validate_file(self, filename: str, file_size: int, max_size_mb: int = 50) -> Tuple[bool, str]:
        """Validate uploaded file"""
        # Check file extension
        file_ext = os.path.splitext(filename)[1].lower()
        if file_ext not in self.supported_types:
            return False, f"Unsupported file type: {file_ext}. Supported types: {', '.join(self.supported_types.keys())}"
        
        # Check file size
        max_size_bytes = max_size_mb * 1024 * 1024
        if file_size > max_size_bytes:
            return False, f"File too large. Maximum size: {max_size_mb}MB"
        
        return True, "File is valid"
    
    def get_uploaded_documents(self, limit: int = 100) -> List[Dict[str, Any]]:
        """Retrieve uploaded documents from Qdrant"""
        try:
            # Search for all documents with source_type = 'user_upload'
            search_results = self.qdrant_client.scroll(
                collection_name=self.collection_name,
                scroll_filter=models.Filter(
                    must=[
                        models.FieldCondition(
                            key="source_type",
                            match=models.MatchValue(value="user_upload")
                        )
                    ]
                ),
                limit=limit,
                with_payload=True
            )
            
            # Group by doc_id to get unique documents
            documents = {}
            for point in search_results[0]:
                doc_id = point.payload.get('doc_id')
                if doc_id and doc_id not in documents:
                    documents[doc_id] = {
                        'doc_id': doc_id,
                        'filename': point.payload.get('filename', 'Unknown'),
                        'title': point.payload.get('title', 'Unknown'),
                        'source': point.payload.get('source', 'User Upload'),
                        'source_type': point.payload.get('source_type', 'user_upload'),
                        'file_type': point.payload.get('file_type', 'unknown'),
                        'file_size': point.payload.get('file_size', 0),
                        'upload_timestamp': point.payload.get('upload_timestamp', ''),
                        'total_chunks': point.payload.get('total_chunks', 0),
                        'categories': ['user_upload']  # Default category for uploaded docs
                    }
            
            return list(documents.values())
            
        except Exception as e:
            self.logger.error(f"Error retrieving uploaded documents: {e}")
            return []
    
    def delete_uploaded_document(self, doc_id: str) -> Dict[str, Any]:
        """Delete an uploaded document from Qdrant"""
        try:
            # Delete all points with the specified doc_id
            self.qdrant_client.delete(
                collection_name=self.collection_name,
                points_selector=models.Filter(
                    must=[
                        models.FieldCondition(
                            key="doc_id",
                            match=models.MatchValue(value=doc_id)
                        ),
                        models.FieldCondition(
                            key="source_type",
                            match=models.MatchValue(value="user_upload")
                        )
                    ]
                )
            )
            
            self.logger.info(f"Successfully deleted document {doc_id}")
            return {
                'success': True,
                'doc_id': doc_id,
                'message': f'Document {doc_id} deleted successfully'
            }
            
        except Exception as e:
            self.logger.error(f"Error deleting document {doc_id}: {e}")
            return {
                'success': False,
                'doc_id': doc_id,
                'error': str(e)
            } 